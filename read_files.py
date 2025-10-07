from pptx import Presentation
import openpyxl

prs = Presentation('Presentation.pptx')
print("=== POWERPOINT CONTENT ===\n")
for i, slide in enumerate(prs.slides, 1):
    print(f"--- Slide {i} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
    print()

wb = openpyxl.load_workbook('Darden Financials.xlsx', data_only=True)
print("\n=== EXCEL CONTENT ===\n")
for sheet_name in wb.sheetnames:
    print(f"--- Sheet: {sheet_name} ---")
    ws = wb[sheet_name]
    for row in ws.iter_rows(values_only=True):
        if any(cell is not None for cell in row):
            print(row)
    print()
