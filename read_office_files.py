import openpyxl
from pptx import Presentation
import zipfile

print("=" * 80)
print("NOTES.ZIP CONTENTS")
print("=" * 80)
with zipfile.ZipFile('notes.zip', 'r') as zip_ref:
    print(zip_ref.namelist())

print("\n" + "=" * 80)
print("EXCEL FILE: Darden Financials.xlsx")
print("=" * 80)
wb = openpyxl.load_workbook('Darden Financials.xlsx', data_only=True)
for sheet_name in wb.sheetnames:
    print(f"\n--- Sheet: {sheet_name} ---")
    ws = wb[sheet_name]
    for row in ws.iter_rows(values_only=True):
        if any(cell is not None for cell in row):
            print(row)

print("\n" + "=" * 80)
print("POWERPOINT FILE: Presentation.pptx")
print("=" * 80)
prs = Presentation('Presentation.pptx')
for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- Slide {i} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text)
